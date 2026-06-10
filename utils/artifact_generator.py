import io
import re
import pandas as pd
import matplotlib.pyplot as plt
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation


def generate_table(content: str) -> pd.DataFrame:
    """Convert a textual table (CSV or markdown) into a pandas DataFrame.

    - Handles comma‑separated rows ("a,b,c")
    - Detects markdown tables using pipe delimiters ("| Col | Value |")
    - If quarterly data spans multiple rows with a year pattern, optionally aggregates to yearly totals.
    """
    if "|" in content:
        lines = [ln.strip() for ln in content.strip().split("\n") if ln.strip()]
        rows = [ln for ln in lines if not re.match(r"^\s*\|?\s*-+\s*\|?", ln)]
        parsed = []
        for row in rows:
            row = row.strip("|")
            cells = [c.strip() for c in row.split("|")]
            parsed.append(cells)
        if not parsed:
            return pd.DataFrame()
        max_cols = max(len(r) for r in parsed)
        padded = [r + [""] * (max_cols - len(r)) for r in parsed]
        header = padded[0]
        data = padded[1:] if len(padded) > 1 else []
        df = pd.DataFrame(data, columns=header)
        if "Quarter" in df.columns and any(df["Quarter"].str.contains("Q")):
            df["Year"] = df["Quarter"].str.extract(r"(\\d{4})")
            sales_col = next(
                (
                    c
                    for c in df.columns
                    if c.lower() in ["sales", "automotive sales", "sales amount"]
                ),
                None,
            )
            if sales_col:
                df[sales_col] = (
                    df[sales_col].replace("[,]", "", regex=True).astype(float)
                )
                yearly = df.groupby("Year")[sales_col].sum().reset_index()
                yearly.columns = ["Year", sales_col]
                return yearly
        return df
    rows = [r.split(",") for r in content.strip().split("\n") if r]
    if not rows:
        return pd.DataFrame()
    col_counts = {len(r) for r in rows}
    if len(col_counts) == 1:
        header = rows[0]
        data = rows[1:]
        return pd.DataFrame(data, columns=header)
    # Fallback single‑column DataFrame
    return pd.DataFrame({"Answer": [content]})


def generate_excel(df: pd.DataFrame) -> bytes:
    with io.BytesIO() as buffer:
        with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
            df.to_excel(writer, index=False)
        return buffer.getvalue()


def generate_chart(content: str, chart_type: str = "line") -> bytes:
    """Generate a chart (line, bar, pie, histogram) from extracted numeric or key‑value data.

    - Parses "label: value" pairs for categorical charts.
    - Falls back to numeric list for line, bar, histogram.
    - chart_type can be "line", "bar", "pie", or "histogram".
    """
    # Extract label‑value pairs from the whole content using a flexible regex
    # Supports formats like "Label: 123", "Label - $123k", "Label = 123%", or "Label 123"
    pairs = []
    pattern = re.compile(
        r"(?P<label>[\w\-\s]+?)\s*[:\-\–=]\s*\$?\s*(?P<value>[-+]?\d[\d,\.]*\s*[kKmM%]?)"
    )
    for match in pattern.finditer(content):
        label = match.group("label").strip()
        # Remove any leading bullet characters that might have been captured
        label = label.lstrip("-*• ").strip()
        raw_val = (
            match.group("value")
            .replace(",", "")
            .replace("$", "")
            .replace("%", "")
            .strip()
        )
        # Handle magnitude suffixes (k, m)
        multiplier = 1
        if raw_val.lower().endswith("k"):
            multiplier = 1_000
            raw_val = raw_val[:-1]
        elif raw_val.lower().endswith("m"):
            multiplier = 1_000_000
            raw_val = raw_val[:-1]
        try:
            value = float(raw_val) * multiplier
            pairs.append((label, value))
        except ValueError:
            continue
    # If no pairs were found, fall back to extracting any numbers in the text
    if not pairs:
        values = [
            float(x.replace(",", ""))
            for x in re.findall(r"[-+]?[0-9]*\.?[0-9]+", content)
        ]
        labels = [str(i) for i in range(len(values))]
    else:
        labels, values = zip(*pairs)

    if pairs:
        labels, values = zip(*pairs)
    else:
        values = [
            float(x.replace(",", ""))
            for x in re.findall(r"[-+]?[0-9]*\.?[0-9]+", content)
        ]
        labels = [str(i) for i in range(len(values))]
    labels = list(labels)
    values = list(values)
    if chart_type == "pie":
        positive = [(lbl, val) for lbl, val in zip(labels, values) if val > 0]
        if positive:
            labels, values = zip(*positive)
            labels = list(labels)
            values = list(values)
        else:
            chart_type = "line"
    plt.figure(figsize=(6, 4))
    if chart_type == "pie" and values:
        plt.pie(values, labels=labels, autopct="%1.1f%%")
        plt.title("Pie Chart")
    elif chart_type == "histogram":
        plt.hist(values, bins="auto", edgecolor="black")
        plt.title("Histogram")
    elif chart_type == "bar":
        plt.bar(labels, values, color="skyblue")
        plt.title("Bar Chart")
    else:
        plt.plot(values, marker="o")
        plt.title("Line Chart")
    buf = io.BytesIO()
    plt.savefig(buf, format="png")
    plt.close()
    return buf.getvalue()


def generate_docx(content: str) -> bytes:
    doc = Document()
    doc.add_heading("Generated Report", level=1)
    doc.add_paragraph(content)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generate_pdf(content: str) -> bytes:

    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    text_obj = c.beginText()
    text_obj.setTextOrigin(40, 750)
    text_obj.setFont("Helvetica", 10)
    for line in content.split("\n"):
        text_obj.textLine(line)
    c.drawText(text_obj)
    c.showPage()
    c.save()
    buffer.seek(0)
    return buffer.getvalue()


def generate_pptx(content: str) -> bytes:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Generated Presentation"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = content
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
