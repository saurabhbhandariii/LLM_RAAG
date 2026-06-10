"""
Privacy-Aware Output Formatter
Generates tables, PDFs, Excel, charts, docs, and presentations with PII sanitization
"""

import io
import re
import logging
from typing import Optional, Dict, List, Tuple
import pandas as pd
import matplotlib.pyplot as plt
from datetime import datetime

# Import privacy utilities
try:
    from guard_utils import PrivacySanitizer
    from privacy_config import get_privacy_config
except ImportError:
    PrivacySanitizer = None
    get_privacy_config = None

logger = logging.getLogger(__name__)


class PrivacyAwareOutputFormatter:
    """Generate various output formats with automatic PII redaction"""

    def __init__(self, sanitize: bool = True):
        self.sanitize = sanitize and PrivacySanitizer is not None

    def _sanitize_if_needed(self, text: str) -> str:
        """Conditionally sanitize text based on privacy settings"""
        if self.sanitize:
            return PrivacySanitizer.sanitize_text(text)
        return text

    def _sanitize_dataframe(self, df: pd.DataFrame) -> pd.DataFrame:
        """Sanitize all text columns in a DataFrame"""
        if not self.sanitize:
            return df

        df_copy = df.copy()
        for col in df_copy.select_dtypes(include=["object"]).columns:
            df_copy[col] = df_copy[col].apply(
                lambda x: self._sanitize_if_needed(str(x)) if x else x
            )
        return df_copy

    # ========================================================================
    # TABLE GENERATION
    # ========================================================================

    def generate_table(self, content: str) -> pd.DataFrame:
        """
        Convert textual table (CSV or markdown) into sanitized pandas DataFrame.

        Features:
        - Markdown table parsing (pipe delimiters)
        - CSV parsing
        - Automatic yearly aggregation for quarterly data
        - PII sanitization on all cells
        """
        # Sanitize input content first
        content = self._sanitize_if_needed(content)

        if "|" in content:
            # Markdown table parsing
            lines = [ln.strip() for ln in content.strip().split("\n") if ln.strip()]
            rows = [
                ln for ln in lines if not re.match(r"^\s*\|?\s*-+\s*\|?", ln)
            ]
            parsed = []
            for row in rows:
                row = row.strip("|")
                cells = [c.strip() for c in row.split("|")]
                parsed.append(cells)

            if not parsed:
                return pd.DataFrame()

            max_cols = max(len(r) for r in parsed)
            padded = [r + [""] * (max_cols - len(r)) for r in parsed]
            header = padded[0]
            data = padded[1:] if len(padded) > 1 else []
            df = pd.DataFrame(data, columns=header)

            # Quarterly aggregation logic
            if "Quarter" in df.columns and any(df["Quarter"].str.contains("Q")):
                df["Year"] = df["Quarter"].str.extract(r"(\d{4})")
                sales_col = next(
                    (
                        c
                        for c in df.columns
                        if c.lower()
                        in ["sales", "automotive sales", "sales amount"]
                    ),
                    None,
                )
                if sales_col:
                    df[sales_col] = (
                        df[sales_col].replace("[,]", "", regex=True).astype(float)
                    )
                    yearly = df.groupby("Year")[sales_col].sum().reset_index()
                    yearly.columns = ["Year", sales_col]
                    df = yearly

        else:
            # CSV parsing
            rows = [r.split(",") for r in content.strip().split("\n") if r]
            if not rows:
                return pd.DataFrame()

            col_counts = {len(r) for r in rows}
            if len(col_counts) == 1:
                header = rows[0]
                data = rows[1:]
                df = pd.DataFrame(data, columns=header)
            else:
                df = pd.DataFrame({"Answer": [content]})

        # Sanitize and return
        return self._sanitize_dataframe(df)

    # ========================================================================
    # EXCEL GENERATION
    # ========================================================================

    def generate_excel(
        self, df: pd.DataFrame, include_metadata: bool = True
    ) -> bytes:
        """
        Generate Excel file with styled tables and metadata.

        Features:
        - PII sanitization
        - Auto-formatted columns
        - Metadata sheet with timestamp and source info
        """
        df = self._sanitize_dataframe(df)

        with io.BytesIO() as buffer:
            with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
                df.to_excel(writer, sheet_name="Data", index=False)

                if include_metadata:
                    metadata_df = pd.DataFrame(
                        {
                            "Property": [
                                "Generated",
                                "Privacy Level",
                                "Sanitized",
                                "Rows",
                                "Columns",
                            ],
                            "Value": [
                                datetime.utcnow().isoformat(),
                                get_privacy_config()
                                .get("level", "STANDARD")
                                if get_privacy_config
                                else "STANDARD",
                                str(self.sanitize),
                                str(len(df)),
                                str(len(df.columns)),
                            ],
                        }
                    )
                    metadata_df.to_excel(
                        writer, sheet_name="Metadata", index=False
                    )

                # Auto-format columns
                for sheet_name, sheet in writer.sheets.items():
                    for column in sheet.columns:
                        max_length = max(
                            (
                                len(str(cell.value))
                                for cell in column
                                if cell.value
                            ),
                            default=10,
                        )
                        sheet.column_dimensions[column[0].column_letter].width = (
                            min(max_length + 2, 50)
                        )

            return buffer.getvalue()

    # ========================================================================
    # CHART GENERATION
    # ========================================================================

    def generate_chart(
        self,
        content: str,
        chart_type: str = "line",
        title: str = "",
        figsize: Tuple[int, int] = (10, 6),
    ) -> bytes:
        """
        Generate chart from text data with styling.

        Features:
        - Automatic label-value extraction
        - Multiple chart types (line, bar, pie, histogram)
        - Professional styling
        - Sanitized labels and values
        """
        # Sanitize input
        content = self._sanitize_if_needed(content)

        # Extract label-value pairs
        pairs = []
        pattern = re.compile(
            r"(?P<label>[\w\-\s]+?)\s*[:\-\–=]\s*\$?\s*(?P<value>[-+]?\d[\d,\.]*\s*[kKmM%]?)"
        )

        for match in pattern.finditer(content):
            label = match.group("label").strip().lstrip("-*• ").strip()
            raw_val = (
                match.group("value")
                .replace(",", "")
                .replace("$", "")
                .replace("%", "")
                .strip()
            )

            # Handle magnitude suffixes
            multiplier = 1
            if raw_val.lower().endswith("k"):
                multiplier = 1_000
                raw_val = raw_val[:-1]
            elif raw_val.lower().endswith("m"):
                multiplier = 1_000_000
                raw_val = raw_val[:-1]

            try:
                value = float(raw_val) * multiplier
                pairs.append((label, value))
            except ValueError:
                continue

        # Fallback if no pairs found
        if not pairs:
            values = [
                float(x.replace(",", ""))
                for x in re.findall(r"[-+]?[0-9]*\.?[0-9]+", content)
            ]
            labels = [f"Value {i}" for i in range(len(values))]
        else:
            labels, values = zip(*pairs)
            labels = list(labels)
            values = list(values)

        # Filter negative values for pie chart
        if chart_type == "pie":
            positive = [(lbl, val) for lbl, val in zip(labels, values) if val > 0]
            if positive:
                labels, values = zip(*positive)
                labels = list(labels)
                values = list(values)
            else:
                chart_type = "line"

        # Create chart
        plt.style.use("seaborn-v0_8-darkgrid")
        plt.figure(figsize=figsize)

        if chart_type == "pie" and values:
            colors = plt.cm.Set3(range(len(values)))
            plt.pie(
                values, labels=labels, autopct="%1.1f%%", colors=colors, startangle=90
            )
            plt.title(title or "Pie Chart", fontsize=14, fontweight="bold")

        elif chart_type == "histogram":
            plt.hist(values, bins="auto", edgecolor="black", color="steelblue")
            plt.title(title or "Histogram", fontsize=14, fontweight="bold")
            plt.xlabel("Value")
            plt.ylabel("Frequency")

        elif chart_type == "bar":
            colors = plt.cm.Set2(range(len(labels)))
            plt.bar(labels, values, color=colors, edgecolor="black")
            plt.title(title or "Bar Chart", fontsize=14, fontweight="bold")
            plt.xlabel("Category")
            plt.ylabel("Value")
            plt.xticks(rotation=45, ha="right")

        else:  # line chart (default)
            plt.plot(values, marker="o", linewidth=2, markersize=8, color="steelblue")
            plt.xticks(range(len(labels)), labels, rotation=45, ha="right")
            plt.title(title or "Line Chart", fontsize=14, fontweight="bold")
            plt.xlabel("Category")
            plt.ylabel("Value")

        plt.tight_layout()

        buf = io.BytesIO()
        plt.savefig(buf, format="png", dpi=300)
        plt.close()

        return buf.getvalue()

    # ========================================================================
    # GRAPH GENERATION
    # ========================================================================

    def generate_graph(
        self,
        nodes: List[str],
        edges: List[Tuple[str, str]],
        title: str = "Network Graph",
    ) -> bytes:
        """
        Generate network graph visualization.

        Features:
        - Node and edge sanitization
        - Force-directed layout
        - Professional styling
        """
        try:
            import networkx as nx
        except ImportError:
            logger.error("networkx not installed for graph generation")
            return self._generate_chart_fallback(
                f"Graph: {title}\nNodes: {len(nodes)}, Edges: {len(edges)}"
            )

        # Sanitize nodes and edges
        nodes = [self._sanitize_if_needed(node) for node in nodes]
        edges = [
            (self._sanitize_if_needed(u), self._sanitize_if_needed(v))
            for u, v in edges
        ]

        # Create graph
        G = nx.Graph()
        G.add_nodes_from(nodes)
        G.add_edges_from(edges)

        plt.figure(figsize=(12, 8))
        pos = nx.spring_layout(G, k=2, iterations=50)

        nx.draw_networkx_nodes(
            G, pos, node_color="lightblue", node_size=1500, alpha=0.8
        )
        nx.draw_networkx_edges(G, pos, alpha=0.5, width=2)
        nx.draw_networkx_labels(G, pos, font_size=8, font_weight="bold")

        plt.title(title, fontsize=14, fontweight="bold")
        plt.axis("off")
        plt.tight_layout()

        buf = io.BytesIO()
        plt.savefig(buf, format="png", dpi=300)
        plt.close()

        return buf.getvalue()

    # ========================================================================
    # DOCUMENT GENERATION
    # ========================================================================

    def generate_docx(self, content: str, title: str = "Generated Report") -> bytes:
        """
        Generate Word document with styling.

        Features:
        - PII sanitization
        - Headings and formatting
        - Metadata footer
        """
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        content = self._sanitize_if_needed(content)

        doc = Document()
        doc.add_heading(title, level=1)

        # Add timestamp
        metadata = doc.add_paragraph()
        metadata.text = (
            f"Generated: {datetime.utcnow().isoformat()} | Sanitized: {self.sanitize}"
        )
        metadata.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        metadata.runs[0].font.size = 100  # 8pt

        doc.add_paragraph()  # Spacing

        # Add content with proper paragraphs
        for line in content.split("\n"):
            if line.strip():
                doc.add_paragraph(line)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.getvalue()

    def generate_pdf(self, content: str, title: str = "Generated Report") -> bytes:
        """
        Generate PDF with formatting.

        Features:
        - PII sanitization
        - Multiple pages support
        - Professional layout
        """
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch

        content = self._sanitize_if_needed(content)

        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter

        # Title
        c.setFont("Helvetica-Bold", 16)
        c.drawString(0.5 * inch, height - 0.5 * inch, title)

        # Metadata
        c.setFont("Helvetica-Oblique", 8)
        c.drawString(
            0.5 * inch,
            height - 0.75 * inch,
            f"Generated: {datetime.utcnow().isoformat()} | Sanitized: {self.sanitize}",
        )

        # Content
        c.setFont("Helvetica", 10)
        y = height - 1.2 * inch
        line_height = 0.2 * inch

        for line in content.split("\n"):
            if y < 0.5 * inch:  # New page
                c.showPage()
                y = height - 0.5 * inch

            c.drawString(0.5 * inch, y, line[:80])  # Line length limit
            y -= line_height

        c.showPage()
        c.save()
        buffer.seek(0)
        return buffer.getvalue()

    def generate_pptx(
        self, content: str, title: str = "Generated Presentation"
    ) -> bytes:
        """
        Generate PowerPoint presentation.

        Features:
        - PII sanitization
        - Multi-slide layout
        - Bullet points
        """
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()
        prs.slide_width = prs.slide_height = None

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        subtitle = slide.placeholders[1].text_frame
        subtitle.text = f"Generated: {datetime.utcnow().isoformat()}"

        # Content slides (split by bullet points or line count)
        content = self._sanitize_if_needed(content)
        content_slide_layout = prs.slide_layouts[1]

        # Split content into chunks
        lines = content.split("\n")
        chunk_size = 10

        for i in range(0, len(lines), chunk_size):
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = f"Content - Part {i // chunk_size + 1}"

            body = slide.placeholders[1].text_frame
            body.clear()

            for line in lines[i : i + chunk_size]:
                if line.strip():
                    p = body.add_paragraph()
                    p.text = line.strip()
                    p.level = 0

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()

    # ========================================================================
    # FALLBACK METHODS
    # ========================================================================

    def _generate_chart_fallback(self, content: str) -> bytes:
        """Fallback if chart generation fails"""
        plt.figure(figsize=(8, 6))
        plt.text(0.5, 0.5, content, ha="center", va="center", fontsize=12)
        plt.axis("off")

        buf = io.BytesIO()
        plt.savefig(buf, format="png", dpi=100)
        plt.close()
        return buf.getvalue()


# ============================================================================
# PUBLIC API (backward compatible)
# ============================================================================

_formatter = PrivacyAwareOutputFormatter(sanitize=True)


def generate_table(content: str) -> pd.DataFrame:
    """Public API: Generate table"""
    return _formatter.generate_table(content)


def generate_excel(df: pd.DataFrame) -> bytes:
    """Public API: Generate Excel"""
    return _formatter.generate_excel(df)


def generate_chart(content: str, chart_type: str = "line") -> bytes:
    """Public API: Generate chart"""
    return _formatter.generate_chart(content, chart_type)


def generate_docx(content: str) -> bytes:
    """Public API: Generate Word document"""
    return _formatter.generate_docx(content)


def generate_pdf(content: str) -> bytes:
    """Public API: Generate PDF"""
    return _formatter.generate_pdf(content)


def generate_pptx(content: str) -> bytes:
    """Public API: Generate PowerPoint"""
    return _formatter.generate_pptx(content)
